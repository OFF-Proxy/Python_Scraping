from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "デュエル開始の宣言をしろ！磯野！"
subtitle.text = "デュエル開始～！"

# スライドごとに文字数をカウント
total_character_count = 0
for slide in prs.slides:
    slide_character_count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                slide_character_count += len(paragraph.text)
    total_character_count += slide_character_count

# 文字数を表示
print("プレゼンテーション内の文字数:", total_character_count)

prs.save('test2.pptx')